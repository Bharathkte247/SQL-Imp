from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


OUTPUT = Path("Report_Gap_Analysis_Presentation.pptx")

NAVY = RGBColor(18, 42, 74)
BLUE = RGBColor(42, 117, 185)
TEAL = RGBColor(0, 150, 136)
GREEN = RGBColor(76, 175, 80)
AMBER = RGBColor(245, 158, 11)
RED = RGBColor(220, 53, 69)
PURPLE = RGBColor(111, 66, 193)
GRAY = RGBColor(91, 103, 112)
LIGHT_GRAY = RGBColor(243, 246, 249)
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(30, 30, 30)


def set_fill(shape, color, transparency=0):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.fill.transparency = transparency
    shape.line.color.rgb = WHITE


def set_text(
    shape,
    text,
    size=18,
    color=BLACK,
    bold=False,
    align=PP_ALIGN.LEFT,
    font="Aptos",
):
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.12)
    frame.margin_right = Inches(0.12)
    frame.margin_top = Inches(0.08)
    frame.margin_bottom = Inches(0.08)
    p = frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return shape


def add_title(slide, title, subtitle=None):
    box = slide.shapes.add_textbox(Inches(0.45), Inches(0.2), Inches(12.4), Inches(0.45))
    set_text(box, title, size=24, color=NAVY, bold=True)
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.48), Inches(0.67), Inches(11.9), Inches(0.28))
        set_text(sub, subtitle, size=10.5, color=GRAY)
    line = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0.48),
        Inches(0.98),
        Inches(12.35),
        Inches(0.03),
    )
    set_fill(line, BLUE)


def add_footer(slide, idx):
    footer = slide.shapes.add_textbox(Inches(0.45), Inches(7.15), Inches(12.4), Inches(0.22))
    set_text(footer, f"Report & Dashboard Gap Analysis | {idx}", size=8.5, color=GRAY, align=PP_ALIGN.RIGHT)


def add_bullets(slide, x, y, w, h, bullets, size=13, color=BLACK):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Aptos"
    return box


def add_card(slide, x, y, w, h, title, body, color=BLUE):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    set_fill(shape, color)
    tf = shape.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = title
    r.font.size = Pt(18)
    r.font.bold = True
    r.font.color.rgb = WHITE
    r.font.name = "Aptos Display"
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = body
    r2.font.size = Pt(10.5)
    r2.font.color.rgb = WHITE
    r2.font.name = "Aptos"
    return shape


def add_table(slide, x, y, w, h, headers, rows, widths=None, font_size=8.8):
    table_shape = slide.shapes.add_table(
        len(rows) + 1,
        len(headers),
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    table = table_shape.table
    if widths:
        for idx, width in enumerate(widths):
            table.columns[idx].width = Inches(width)

    for col, header in enumerate(headers):
        cell = table.cell(0, col)
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        cell.text = header
        para = cell.text_frame.paragraphs[0]
        para.font.bold = True
        para.font.size = Pt(font_size)
        para.font.color.rgb = WHITE
        para.font.name = "Aptos"

    for r_idx, row in enumerate(rows, start=1):
        for c_idx, value in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_GRAY if r_idx % 2 else WHITE
            cell.text = str(value)
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(font_size)
            para.font.color.rgb = BLACK
            para.font.name = "Aptos"
    return table_shape


def add_chevrons(slide, x, y, w, h, labels, colors):
    step_w = w / len(labels)
    for idx, label in enumerate(labels):
        shape_type = MSO_AUTO_SHAPE_TYPE.CHEVRON if idx < len(labels) - 1 else MSO_AUTO_SHAPE_TYPE.PENTAGON
        shape = slide.shapes.add_shape(
            shape_type,
            Inches(x + idx * step_w),
            Inches(y),
            Inches(step_w + 0.08),
            Inches(h),
        )
        set_fill(shape, colors[idx])
        set_text(shape, label, size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)


def add_connector_arrow(slide, x, y, w, color=GRAY):
    arrow = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(0.32),
    )
    set_fill(arrow, color)
    return arrow


def build_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # Slide 1
    slide = prs.slides.add_slide(blank)
    bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    set_fill(bg, NAVY)
    accent = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(6.85), prs.slide_width, Inches(0.65))
    set_fill(accent, BLUE)
    title = slide.shapes.add_textbox(Inches(0.75), Inches(1.05), Inches(11.8), Inches(0.9))
    set_text(title, "Report & Dashboard Gap Analysis", size=36, color=WHITE, bold=True)
    subtitle = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(10.8), Inches(0.55))
    set_text(subtitle, "Deployment Impact Assessment for Standard Reports, Dashboards, and Reporting Fields", size=18, color=WHITE)
    add_chevrons(
        slide,
        0.95,
        3.25,
        11.3,
        0.8,
        ["Deployment Changes", "Report Inventory", "Gap Analysis", "Remediation Plan"],
        [BLUE, TEAL, AMBER, GREEN],
    )
    note = slide.shapes.add_textbox(Inches(0.85), Inches(5.25), Inches(11.4), Inches(0.65))
    set_text(note, "Prepared for business, product, data, and reporting stakeholders", size=16, color=WHITE, align=PP_ALIGN.CENTER)

    # Slide 2
    slide = prs.slides.add_slide(blank)
    add_title(slide, "Executive Summary", "High-level summary of what the latest deployment changes mean for reporting continuity.")
    add_card(slide, 0.65, 1.35, 3.7, 1.2, "What Changed", "Fields, values, calculations, and dashboard dependencies", BLUE)
    add_card(slide, 4.8, 1.35, 3.7, 1.2, "What May Break", "Filters, KPIs, historical comparisons, and visuals", AMBER)
    add_card(slide, 8.95, 1.35, 3.7, 1.2, "What We Need", "Updates, validation, sign-off, and communication", TEAL)
    add_bullets(
        slide,
        0.9,
        3.05,
        5.6,
        2.6,
        [
            "Prioritize high-usage and executive-facing reports first.",
            "Assess field-level dependencies before changing dashboard visuals.",
            "Validate metric deltas with business owners to avoid false defects.",
            "Publish a clear change log for report consumers.",
        ],
        size=14,
    )
    callout = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(7.05), Inches(3.05), Inches(5.45), Inches(2.3))
    set_fill(callout, LIGHT_GRAY)
    callout.line.color.rgb = BLUE
    set_text(
        callout,
        "Decision Required\nConfirm the critical report list, field mapping ownership, and approval path before publishing updated reporting assets.",
        size=15,
        color=NAVY,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_footer(slide, 2)

    # Slide 3
    slide = prs.slides.add_slide(blank)
    add_title(slide, "Analysis Approach", "SmartArt-style process used to identify reporting gaps and required remediation.")
    add_chevrons(
        slide,
        0.75,
        1.45,
        11.85,
        1.0,
        ["1. Capture Changes", "2. Map Dependencies", "3. Score Impact", "4. Remediate", "5. Validate & Sign Off"],
        [BLUE, TEAL, PURPLE, AMBER, GREEN],
    )
    rows = [
        ("Capture Changes", "Deployment notes, field changes, value changes, logic updates"),
        ("Map Dependencies", "Reports, dashboards, datasets, filters, formulas, calculations"),
        ("Score Impact", "Criticality, usage frequency, user segment, compliance exposure"),
        ("Remediate", "Update report fields, filters, dashboard tiles, formulas, security"),
        ("Validate", "Compare outputs, test refresh, obtain business owner approval"),
    ]
    add_table(slide, 0.95, 3.1, 11.4, 2.75, ["Step", "Primary Activities"], rows, widths=[2.6, 8.8], font_size=11)
    add_footer(slide, 3)

    # Slide 4
    slide = prs.slides.add_slide(blank)
    add_title(slide, "Deployment-to-Reporting Impact Flow", "How upstream deployment changes can cascade into reports and dashboards.")
    flow_items = [
        ("Deployment Change", "New or modified fields, values, business rules", BLUE),
        ("Data Model / Dataset", "Mapping, refresh, transformation, history", TEAL),
        ("Report Logic", "Filters, formulas, grouping, calculations", PURPLE),
        ("Dashboard Visuals", "Cards, charts, slicers, drill-through views", AMBER),
        ("Business Decision", "KPI interpretation, operations, compliance", GREEN),
    ]
    x_positions = [0.55, 3.15, 5.75, 8.35, 10.95]
    for idx, ((title, body, color), x) in enumerate(zip(flow_items, x_positions)):
        add_card(slide, x, 2.0, 1.9, 1.45, title, body, color)
        if idx < len(flow_items) - 1:
            add_connector_arrow(slide, x + 1.95, 2.55, 0.65, GRAY)
    add_bullets(
        slide,
        1.05,
        4.35,
        11.35,
        1.55,
        [
            "A small field or value-list change can have a visible dashboard impact if it is used in filters or KPI formulas.",
            "Impact assessment should trace from source change to dataset, report logic, dashboard visual, and business outcome.",
        ],
        size=14,
        color=NAVY,
    )
    add_footer(slide, 4)

    # Slide 5
    slide = prs.slides.add_slide(blank)
    add_title(slide, "Report Inventory & Criticality", "Starting inventory to confirm which reporting assets need deeper analysis.")
    rows = [
        ("Executive Dashboard", "Leadership", "Daily", "High", "Yes"),
        ("Sales Performance Report", "Sales Ops", "Daily", "High", "Yes"),
        ("Customer Pipeline Report", "CRM Team", "Weekly", "Medium", "Yes"),
        ("Operational Summary", "Operations", "Monthly", "Low", "No"),
        ("Compliance Extract", "Risk / Compliance", "Monthly", "High", "TBD"),
    ]
    add_table(
        slide,
        0.55,
        1.3,
        12.2,
        2.25,
        ["Report / Dashboard", "Owner", "Frequency", "Criticality", "Impacted?"],
        rows,
        widths=[3.2, 2.2, 1.6, 1.6, 1.7],
        font_size=9.5,
    )
    add_chevrons(
        slide,
        1.15,
        4.35,
        10.95,
        0.8,
        ["Inventory", "Criticality", "Dependency Mapping", "Impact Confirmation"],
        [BLUE, PURPLE, AMBER, GREEN],
    )
    add_footer(slide, 5)

    # Slide 6
    slide = prs.slides.add_slide(blank)
    add_title(slide, "Gap Analysis Heat Map", "Severity scoring across representative reports and deployment change categories.")
    headers = ["Report / Dashboard", "New Fields", "Modified Logic", "Deprecated Fields", "Filters / Values", "Overall"]
    rows = [
        ("Executive Dashboard", "Medium", "High", "Low", "High", "High"),
        ("Sales Performance", "High", "High", "Medium", "Medium", "High"),
        ("Customer Pipeline", "Medium", "Low", "Low", "High", "Medium"),
        ("Compliance Extract", "Low", "Medium", "High", "Medium", "High"),
    ]
    table_shape = add_table(slide, 0.45, 1.35, 12.45, 3.0, headers, rows, widths=[2.7, 1.6, 1.8, 1.9, 1.8, 1.5], font_size=8.6)
    color_map = {"High": RED, "Medium": AMBER, "Low": GREEN}
    table = table_shape.table
    for r_idx in range(1, len(rows) + 1):
        for c_idx in range(1, len(headers)):
            value = table.cell(r_idx, c_idx).text
            table.cell(r_idx, c_idx).fill.solid()
            table.cell(r_idx, c_idx).fill.fore_color.rgb = color_map.get(value, LIGHT_GRAY)
            para = table.cell(r_idx, c_idx).text_frame.paragraphs[0]
            para.font.bold = True
            para.font.color.rgb = WHITE
            para.alignment = PP_ALIGN.CENTER
    legend_items = [("High", RED), ("Medium", AMBER), ("Low", GREEN)]
    for idx, (label, color) in enumerate(legend_items):
        box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(4.1 + idx * 1.75), Inches(5.05), Inches(1.25), Inches(0.45))
        set_fill(box, color)
        set_text(box, label, size=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, 6)

    # Slide 7
    slide = prs.slides.add_slide(blank)
    add_title(slide, "Field-Level Impact Mapping", "Field changes should be mapped to exact reports, formulas, filters, and downstream users.")
    rows = [
        ("Status", "Filter + chart series", "New status values added", "Incomplete views", "Update filters and legends"),
        ("Revenue Amount", "KPI + trend", "Calculation logic updated", "Metric variance", "Validate formula and annotate variance"),
        ("Customer Type", "Grouping", "Deprecated; replaced by Segment", "Blank/error risk", "Replace with Segment"),
        ("Region", "Grouping + security", "No change", "No impact", "No action"),
    ]
    add_table(
        slide,
        0.45,
        1.25,
        12.45,
        2.55,
        ["Field", "Current Usage", "Deployment Change", "Impact", "Recommendation"],
        rows,
        widths=[1.5, 2.25, 2.65, 2.0, 3.0],
        font_size=8.5,
    )
    add_chevrons(
        slide,
        1.0,
        4.6,
        11.2,
        0.75,
        ["Source Field", "Dataset Mapping", "Report Formula", "Dashboard Visual", "User Decision"],
        [BLUE, TEAL, PURPLE, AMBER, GREEN],
    )
    add_footer(slide, 7)

    # Slide 8
    slide = prs.slides.add_slide(blank)
    add_title(slide, "Dashboard Component Impact", "Hub-and-spoke view of dashboard elements that require review after deployment.")
    center = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(5.15), Inches(2.15), Inches(3.0), Inches(1.4))
    set_fill(center, NAVY)
    set_text(center, "Dashboard\nImpact", size=22, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    spokes = [
        ("KPI Cards", 1.0, 1.45, BLUE),
        ("Charts", 4.85, 0.95, TEAL),
        ("Filters", 9.15, 1.45, PURPLE),
        ("Drill-downs", 1.0, 4.3, AMBER),
        ("Refresh Jobs", 4.85, 5.05, RED),
        ("Security / Access", 9.15, 4.3, GREEN),
    ]
    for label, x, y, color in spokes:
        box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(2.25), Inches(0.72))
        set_fill(box, color)
        set_text(box, label, size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, 8)

    # Slide 9
    slide = prs.slides.add_slide(blank)
    add_title(slide, "Business Impact Assessment", "Impact should be explained in business terms, not only technical report changes.")
    tiers = [
        ("High", "Executive decisions, compliance submissions, external reporting", RED, 1.1),
        ("Medium", "Operational tracking, forecast accuracy, adoption confidence", AMBER, 1.55),
        ("Low", "Low-frequency reports or non-critical supporting views", GREEN, 2.0),
    ]
    for idx, (label, desc, color, width_offset) in enumerate(tiers):
        shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.TRAPEZOID,
            Inches(2.2 + width_offset / 2),
            Inches(1.55 + idx * 1.15),
            Inches(8.9 - width_offset),
            Inches(0.95),
        )
        set_fill(shape, color)
        set_text(shape, f"{label} Impact — {desc}", size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_bullets(
        slide,
        1.15,
        5.25,
        11.1,
        1.0,
        [
            "Use impact level to drive remediation order and business sign-off sequence.",
            "Document known metric changes separately from defects so stakeholders understand expected variance.",
        ],
        size=13.5,
        color=NAVY,
    )
    add_footer(slide, 9)

    # Slide 10
    slide = prs.slides.add_slide(blank)
    add_title(slide, "Recommended Remediation Roadmap", "Sequenced actions to update reporting assets and reduce deployment risk.")
    add_chevrons(
        slide,
        0.7,
        1.55,
        11.9,
        0.9,
        ["Prioritize", "Update", "Validate", "Approve", "Publish", "Monitor"],
        [BLUE, TEAL, PURPLE, AMBER, GREEN, NAVY],
    )
    rows = [
        ("Prioritize", "Confirm critical reports, owners, and usage frequency"),
        ("Update", "Replace deprecated fields, adjust formulas, add new values"),
        ("Validate", "Compare metrics, refresh dashboards, test filters and drill-downs"),
        ("Approve", "Obtain business owner sign-off for high-impact reports"),
        ("Publish", "Promote updated assets and communicate report changes"),
        ("Monitor", "Track refresh errors, user feedback, and KPI anomalies"),
    ]
    add_table(slide, 0.95, 3.0, 11.4, 2.85, ["Phase", "Actions"], rows, widths=[2.0, 9.4], font_size=10.2)
    add_footer(slide, 10)

    # Slide 11
    slide = prs.slides.add_slide(blank)
    add_title(slide, "Validation & Sign-Off Checklist", "Checklist to confirm reports are accurate and trusted before rollout.")
    checklist = [
        ("Report runs without errors", GREEN),
        ("Field mappings verified", GREEN),
        ("Filters include new deployment values", AMBER),
        ("KPI formulas reconciled", AMBER),
        ("Dashboard refresh tested", GREEN),
        ("Business owner sign-off captured", BLUE),
    ]
    for idx, (item, color) in enumerate(checklist):
        x = 0.95 + (idx % 2) * 6.05
        y = 1.45 + (idx // 2) * 1.2
        box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(5.45), Inches(0.78))
        set_fill(box, color)
        set_text(box, f"✓ {item}", size=15, color=WHITE, bold=True)
    callout = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(1.05), Inches(5.35), Inches(11.15), Inches(0.75))
    set_fill(callout, LIGHT_GRAY)
    callout.line.color.rgb = BLUE
    set_text(callout, "Recommended evidence: before/after screenshots, data extracts, field mapping notes, and owner approval record.", size=14, color=NAVY, align=PP_ALIGN.CENTER)
    add_footer(slide, 11)

    # Slide 12
    slide = prs.slides.add_slide(blank)
    add_title(slide, "Stakeholder Ownership Model", "Clear ownership prevents gaps from being missed during deployment readiness.")
    rows = [
        ("Product / Deployment Team", "Provide change notes, field definitions, and expected metric behavior"),
        ("Data / BI Team", "Map dependencies, update datasets, modify reports and dashboards"),
        ("Business Owners", "Validate outputs, confirm acceptance, approve report publication"),
        ("Support / Operations", "Monitor post-release issues and user feedback"),
    ]
    add_table(slide, 0.8, 1.35, 11.75, 2.4, ["Stakeholder", "Ownership"], rows, widths=[3.1, 8.65], font_size=10.5)
    add_chevrons(
        slide,
        1.35,
        4.65,
        10.6,
        0.78,
        ["Change Notice", "Impact Review", "Business Validation", "Publication", "Post-Release Monitoring"],
        [BLUE, TEAL, PURPLE, GREEN, NAVY],
    )
    add_footer(slide, 12)

    # Slide 13
    slide = prs.slides.add_slide(blank)
    add_title(slide, "Key Takeaways & Next Steps", "Actions to finalize the report gap analysis and prepare the impacted assets.")
    add_card(slide, 0.8, 1.35, 3.7, 1.25, "1", "Confirm impacted report and dashboard inventory", BLUE)
    add_card(slide, 4.85, 1.35, 3.7, 1.25, "2", "Finalize field-level dependency mapping", TEAL)
    add_card(slide, 8.9, 1.35, 3.7, 1.25, "3", "Prioritize high-impact remediation", AMBER)
    add_bullets(
        slide,
        1.0,
        3.35,
        11.2,
        2.0,
        [
            "Use the heat map to focus effort on high-criticality reports first.",
            "Validate expected KPI variance with business users before go-live communication.",
            "Track remaining open gaps with owner, action, priority, and target release status.",
        ],
        size=15,
        color=NAVY,
    )
    closing = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(2.2), Inches(5.85), Inches(8.95), Inches(0.65))
    set_fill(closing, NAVY)
    set_text(closing, "Outcome: trusted reporting continuity after deployment", size=17, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, 13)

    prs.save(OUTPUT)


if __name__ == "__main__":
    build_deck()
    print(f"Created {OUTPUT}")
